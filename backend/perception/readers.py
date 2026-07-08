"""Leitores por formato de arquivo.

Cada função recebe um Path e devolve um Document normalizado. Separado do
DocumentReader para manter ambos os arquivos enxutos (≤120 linhas).
"""
from __future__ import annotations

import csv
import io
import json
import re
from pathlib import Path

from backend.perception.models import Document, Table


def read_pdf(path: Path) -> Document:
    """Extrai texto e tabelas de um PDF via PyMuPDF."""
    import fitz  # pymupdf

    parts: list[str] = []
    tables: list[Table] = []
    with fitz.open(path) as doc:
        for page in doc:
            parts.append(page.get_text())
            for tbl in page.find_tables():
                tables.append(
                    Table([[c or "" for c in row] for row in tbl.extract()])
                )
        meta = {"pages": doc.page_count}
    return Document(str(path), "pdf", "\n".join(parts), tables, meta)


def read_docx(path: Path) -> Document:
    """Extrai parágrafos e tabelas de um DOCX."""
    import docx

    d = docx.Document(str(path))
    text = "\n".join(p.text for p in d.paragraphs)
    tables = [
        Table([[c.text for c in row.cells] for row in t.rows])
        for t in d.tables
    ]
    return Document(str(path), "docx", text, tables)


def read_xlsx(path: Path) -> Document:
    """Extrai planilhas de um XLSX como tabelas."""
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    tables: list[Table] = []
    texts: list[str] = []
    for ws in wb.worksheets:
        rows = [
            [str(c) if c is not None else "" for c in row]
            for row in ws.iter_rows(values_only=True)
        ]
        if rows:
            tables.append(Table(rows))
            texts.append(f"[{ws.title}]")
    wb.close()
    return Document(str(path), "xlsx", "\n".join(texts), tables)


def read_pptx(path: Path) -> Document:
    """Extrai o texto dos slides de um PPTX."""
    from pptx import Presentation

    prs = Presentation(str(path))
    parts = [
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    ]
    return Document(
        str(path), "pptx", "\n".join(parts),
        metadata={"slides": len(prs.slides._sldIdLst)},
    )


def read_csv(path: Path) -> Document:
    """Lê um CSV como texto e tabela."""
    rows = list(csv.reader(io.StringIO(path.read_text("utf-8"))))
    text = "\n".join(", ".join(r) for r in rows)
    return Document(str(path), "csv", text, [Table(rows)])


def read_json(path: Path) -> Document:
    """Lê um JSON, preservando a estrutura formatada."""
    data = json.loads(path.read_text("utf-8"))
    return Document(
        str(path), "json",
        json.dumps(data, ensure_ascii=False, indent=2),
        metadata={"type": type(data).__name__},
    )


def read_markup(path: Path, fmt: str) -> Document:
    """Remove tags de HTML/XML, deixando o texto limpo."""
    raw = path.read_text("utf-8")
    text = re.sub(r"<[^>]+>", " ", raw)
    text = re.sub(r"\s+", " ", text).strip()
    return Document(str(path), fmt, text)
